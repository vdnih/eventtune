"""Integration ルーター（/api/integration）の統合テスト。

アップロード→バッチ作成→バックグラウンド処理→状態取得のフローを検証する。
AI パイプライン（process_batch / understand_batch）は LLM を呼ぶためフェイクに差し替え、
ルーターの責務（プラン契約・ジョブドキュメントのライフサイクル・stale sweep）に焦点を当てる。
パイプライン内部のロジックは tests/unit/test_ingestion_pipeline.py が担う。
"""

import json
from datetime import UTC, datetime, timedelta

import pytest

pytestmark = pytest.mark.integration


_APPROVED_PLAN = {
    "default_event": {"name": "展示会X", "is_existing": False, "evidence": "ヒントより"},
    "files": [
        {
            "filename": "attendees.csv",
            "business_context": "展示会Xの参加者",
            "targets": [
                {
                    "entity_type": "event_attendances",
                    "column_map": {"name": "name", "email": "email"},
                    "column_modes": {},
                    "link_columns": {},
                }
            ],
            "unmapped_notes": "",
        }
    ],
}


@pytest.fixture
def fake_process_batch(monkeypatch):
    """process_batch を決定的なフェイクへ差し替える。

    ルーターは関数内 import（from agents.data_integration_agent import process_batch）
    のため、モジュール属性の差し替えが呼び出し時に反映される。
    """
    import agents.data_integration_agent as agent
    from agents.data_integration_agent import BatchResult

    calls: list[dict] = []

    async def _fake(files, batch_id, db, plan, space=None):
        calls.append(
            {
                "files": [f for f, _ in files],
                "plan": plan.model_dump() if plan else None,
            }
        )
        return BatchResult(
            created_entities={"persons": 2, "events": 1},
            pending_count=1,
            skipped_count=0,
            report_markdown="# 取り込み結果\n保留 1 件",
        )

    monkeypatch.setattr(agent, "process_batch", _fake)
    return calls


def test_upload_executes_approved_plan_verbatim(make_client, seeded_space, db, fake_process_batch):
    """承認済み BatchPlan がそのまま process_batch に渡る（承認と実行の一致。ADR-015）。"""
    client = make_client(uid="uid_member")
    res = client.post(
        "/api/integration/batches",
        headers={"X-Space-Id": seeded_space},
        files=[("files", ("attendees.csv", b"name,email\nA,a@example.com\n", "text/csv"))],
        data={
            "plan": json.dumps(_APPROVED_PLAN),
            "hint": "展示会Xの参加者リスト",
            "thread_id": "thread_1",
        },
    )
    assert res.status_code == 202
    body_json = res.json()
    batch_id = body_json["batch_id"]
    assert body_json["thread_id"] == "thread_1"

    # TestClient は BackgroundTasks をレスポンス後に同期実行する
    assert len(fake_process_batch) == 1
    executed = fake_process_batch[0]
    assert executed["files"] == ["attendees.csv"]
    assert executed["plan"]["default_event"]["name"] == "展示会X"
    assert executed["plan"]["files"][0]["targets"][0]["column_map"] == {
        "name": "name",
        "email": "email",
    }

    status = client.get(
        f"/api/integration/batches/{batch_id}", headers={"X-Space-Id": seeded_space}
    )
    assert status.status_code == 200
    body = status.json()
    assert body["status"] == "done"
    assert body["created_entities"] == {"persons": 2, "events": 1}
    assert body["pending_count"] == 1
    assert body["report_markdown"].startswith("# 取り込み結果")

    # ジョブドキュメントはスペース配下に作成され、実行されたプランを保持する
    job = db.document(f"spaces/{seeded_space}/integration_jobs/{batch_id}").get()
    assert job.exists
    assert job.to_dict()["plan"]["default_event"]["name"] == "展示会X"

    # 取り込みも会話スレッドと同じ左メニューから見えるよう thread として永続化される
    # （thread_id はフロント採番。batch_id とは別物）
    thread = db.document(f"spaces/{seeded_space}/threads/thread_1").get()
    assert thread.exists
    thread_data = thread.to_dict()
    assert thread_data["uid"] == "uid_member"
    # タイトルはファイル名でなく、承認済みプランの既定イベント名になる
    assert thread_data["title"] == "展示会X"
    docs = list(db.collection(f"spaces/{seeded_space}/threads/thread_1/messages").stream())
    messages = {d.to_dict().get("content_type"): d.to_dict() for d in docs}
    assert messages[None]["role"] == "user"
    assert messages["ingestion_plan"]["batch_id"] == batch_id
    assert messages["ingestion_plan"]["plan"]["default_event"]["name"] == "展示会X"
    assert messages["ingestion_result"]["batch_id"] == batch_id
    assert messages["ingestion_result"]["created_entities"] == {"persons": 2, "events": 1}

    # 承認発話（ボタン起点なら既定文言「取り込む」）が、添付→プラン提示の直後に
    # ユーザー発言として記録される（ボタンとテキスト入力を同じ意味として扱うための土台）。
    assert messages["ingestion_confirm"]["role"] == "user"
    assert messages["ingestion_confirm"]["content"] == "取り込む"
    assert messages["ingestion_confirm"]["batch_id"] == batch_id
    by_seq = sorted(docs, key=lambda d: d.to_dict()["seq"])
    assert [d.to_dict().get("content_type") for d in by_seq] == [
        None,
        "ingestion_plan",
        "ingestion_confirm",
        "ingestion_result",
    ]


def test_upload_persists_custom_confirm_text(make_client, seeded_space, db, fake_process_batch):
    """confirm_text はテキスト入力による承認（ボタン以外）の実際の発言をそのまま保持する。"""
    client = make_client(uid="uid_member")
    res = client.post(
        "/api/integration/batches",
        headers={"X-Space-Id": seeded_space},
        files=[("files", ("attendees.csv", b"name,email\nA,a@example.com\n", "text/csv"))],
        data={
            "plan": json.dumps(_APPROVED_PLAN),
            "thread_id": "thread_confirm_text",
            "confirm_text": "OK、お願いします",
        },
    )
    assert res.status_code == 202
    messages = {
        m.to_dict().get("content_type"): m.to_dict()
        for m in db.collection(
            f"spaces/{seeded_space}/threads/thread_confirm_text/messages"
        ).stream()
    }
    assert messages["ingestion_confirm"]["content"] == "OK、お願いします"


def test_thread_title_uses_dataset_kind_when_no_default_event(
    make_client, seeded_space, db, fake_process_batch
):
    """イベントに紐づかない取り込み（コンテンツリスト等）はスレッドタイトルが種別名になる。"""
    plan = {
        "default_event": None,
        "files": [
            {
                "filename": "contents.csv",
                "business_context": "配布資料一覧",
                "targets": [
                    {
                        "entity_type": "contents",
                        "column_map": {"content_name": "content_name"},
                        "column_modes": {},
                        "link_columns": {},
                    }
                ],
                "unmapped_notes": "",
            }
        ],
    }
    client = make_client(uid="uid_member")
    res = client.post(
        "/api/integration/batches",
        headers={"X-Space-Id": seeded_space},
        files=[("files", ("contents.csv", b"content_name\nContentA\n", "text/csv"))],
        data={"plan": json.dumps(plan), "thread_id": "thread_contents"},
    )
    assert res.status_code == 202

    thread = db.document(f"spaces/{seeded_space}/threads/thread_contents").get()
    assert thread.exists
    assert thread.to_dict()["title"] == "コンテンツリストの取り込み"


def test_plan_omitted_runs_understand_once(make_client, seeded_space, monkeypatch):
    """plan 省略時（API 直叩き）は実行内で Understand が1回だけ走る。"""
    import agents.data_integration_agent as agent
    from agents.data_integration_agent import BatchResult
    from ontology import BatchPlan, FilePlan

    understand_calls: list[list[str]] = []
    executed_plans: list[dict] = []

    async def _fake_understand(files, hint, existing_event_names, space=None):
        understand_calls.append([f for f, _ in files])
        return BatchPlan(files=[FilePlan(filename="a.csv")])

    async def _fake_process(files, batch_id, db, plan, space=None):
        executed_plans.append(plan.model_dump())
        return BatchResult()

    monkeypatch.setattr(agent, "understand_batch", _fake_understand)
    monkeypatch.setattr(agent, "process_batch", _fake_process)

    client = make_client(uid="uid_member")
    res = client.post(
        "/api/integration/batches",
        headers={"X-Space-Id": seeded_space},
        files=[("files", ("a.csv", b"x\n1\n", "text/csv"))],
        data={"thread_id": "thread_2"},
    )
    assert res.status_code == 202
    assert understand_calls == [["a.csv"]]
    assert executed_plans[0]["files"][0]["filename"] == "a.csv"


def test_invalid_plan_json_returns_400(make_client, seeded_space, fake_process_batch):
    client = make_client(uid="uid_member")
    res = client.post(
        "/api/integration/batches",
        headers={"X-Space-Id": seeded_space},
        files=[("files", ("a.csv", b"x\n1\n", "text/csv"))],
        data={"plan": "{broken json", "thread_id": "thread_invalid_plan"},
    )
    assert res.status_code == 400
    assert fake_process_batch == []


def test_legacy_doc_upload_returns_400(make_client, seeded_space, fake_process_batch):
    """旧形式 .doc は明示拒否する（文字化けを AI に渡さない）。"""
    client = make_client(uid="uid_member")
    for endpoint in ("/api/integration/plan", "/api/integration/batches"):
        res = client.post(
            endpoint,
            headers={"X-Space-Id": seeded_space},
            files=[("files", ("report.doc", b"legacy doc bytes", "application/msword"))],
            data={"thread_id": "thread_legacy_doc"},
        )
        assert res.status_code == 400
        assert "未対応のファイル形式" in res.json()["detail"]
    assert fake_process_batch == []


def test_pdf_upload_accepted(make_client, seeded_space, fake_process_batch, monkeypatch):
    """PDF はテキスト抽出限定で受理される（ADR-015 決定7 改訂）。"""
    import agents.data_integration_agent as agent
    from ontology import BatchPlan, FilePlan

    pdf_bytes = b"%PDF-1.4 ..."

    async def _fake_understand(files, hint, existing_event_names, space=None):
        return BatchPlan(files=[FilePlan(filename="report.pdf")])

    monkeypatch.setattr(agent, "understand_batch", _fake_understand)

    client = make_client(uid="uid_member")

    plan_res = client.post(
        "/api/integration/plan",
        headers={"X-Space-Id": seeded_space},
        files=[("files", ("report.pdf", pdf_bytes, "application/pdf"))],
    )
    assert plan_res.status_code == 200
    assert plan_res.json()["files"][0]["extraction_caveat"] != ""

    batch_res = client.post(
        "/api/integration/batches",
        headers={"X-Space-Id": seeded_space},
        files=[("files", ("report.pdf", pdf_bytes, "application/pdf"))],
        data={"thread_id": "thread_pdf"},
    )
    assert batch_res.status_code == 202
    assert fake_process_batch[0]["files"] == ["report.pdf"]


def test_pptx_upload_accepted(make_client, seeded_space, fake_process_batch, monkeypatch):
    """PowerPoint (.pptx) はテキスト抽出限定で受理される。"""
    import io

    from pptx import Presentation

    import agents.data_integration_agent as agent
    from ontology import BatchPlan, FilePlan

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(0, 0, 100, 100)
    box.text_frame.text = "概要"
    buf = io.BytesIO()
    presentation.save(buf)
    pptx_bytes = buf.getvalue()

    async def _fake_understand(files, hint, existing_event_names, space=None):
        return BatchPlan(files=[FilePlan(filename="slides.pptx")])

    monkeypatch.setattr(agent, "understand_batch", _fake_understand)

    client = make_client(uid="uid_member")

    plan_res = client.post(
        "/api/integration/plan",
        headers={"X-Space-Id": seeded_space},
        files=[
            (
                "files",
                (
                    "slides.pptx",
                    pptx_bytes,
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                ),
            )
        ],
    )
    assert plan_res.status_code == 200
    assert plan_res.json()["files"][0]["extraction_caveat"] != ""

    batch_res = client.post(
        "/api/integration/batches",
        headers={"X-Space-Id": seeded_space},
        files=[
            (
                "files",
                (
                    "slides.pptx",
                    pptx_bytes,
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                ),
            )
        ],
        data={"thread_id": "thread_pptx"},
    )
    assert batch_res.status_code == 202
    assert fake_process_batch[0]["files"] == ["slides.pptx"]


def test_docx_upload_accepted(make_client, seeded_space, fake_process_batch, monkeypatch):
    """Word (.docx) は文書として受理される（PDF とは異なり明示拒否しない）。"""
    import io

    import docx as docx_lib

    import agents.data_integration_agent as agent
    from ontology import BatchPlan, FilePlan

    document = docx_lib.Document()
    document.add_paragraph("概要テキスト")
    buf = io.BytesIO()
    document.save(buf)
    docx_bytes = buf.getvalue()

    async def _fake_understand(files, hint, existing_event_names, space=None):
        return BatchPlan(files=[FilePlan(filename="overview.docx")])

    monkeypatch.setattr(agent, "understand_batch", _fake_understand)

    client = make_client(uid="uid_member")

    plan_res = client.post(
        "/api/integration/plan",
        headers={"X-Space-Id": seeded_space},
        files=[
            (
                "files",
                (
                    "overview.docx",
                    docx_bytes,
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                ),
            )
        ],
    )
    assert plan_res.status_code == 200

    batch_res = client.post(
        "/api/integration/batches",
        headers={"X-Space-Id": seeded_space},
        files=[
            (
                "files",
                (
                    "overview.docx",
                    docx_bytes,
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                ),
            )
        ],
        data={"thread_id": "thread_docx"},
    )
    assert batch_res.status_code == 202
    assert fake_process_batch[0]["files"] == ["overview.docx"]


def test_empty_upload_returns_400(make_client, seeded_space, fake_process_batch):
    client = make_client(uid="uid_member")
    res = client.post(
        "/api/integration/batches",
        headers={"X-Space-Id": seeded_space},
        files=[("files", ("empty.csv", b"", "text/csv"))],
        data={"thread_id": "thread_empty"},
    )
    assert res.status_code == 400


def test_pipeline_failure_marks_batch_error(make_client, seeded_space, db, monkeypatch):
    """process_batch が例外を投げてもジョブは error 状態で着地する（握りつぶさない）。"""
    import agents.data_integration_agent as agent

    async def _boom(*args, **kwargs):
        raise RuntimeError("pipeline exploded")

    monkeypatch.setattr(agent, "process_batch", _boom)
    client = make_client(uid="uid_member")
    res = client.post(
        "/api/integration/batches",
        headers={"X-Space-Id": seeded_space},
        files=[("files", ("a.csv", b"x\n1\n", "text/csv"))],
        data={"plan": json.dumps(_APPROVED_PLAN), "thread_id": "thread_failure"},
    )
    batch_id = res.json()["batch_id"]
    body = client.get(
        f"/api/integration/batches/{batch_id}", headers={"X-Space-Id": seeded_space}
    ).json()
    assert body["status"] == "error"
    assert "pipeline exploded" in body["error"]

    # スレッド側にも失敗が記録され、開いたときに分かる
    messages = {
        m.to_dict().get("content_type"): m.to_dict()
        for m in db.collection(f"spaces/{seeded_space}/threads/thread_failure/messages").stream()
    }
    assert "pipeline exploded" in messages["ingestion_error"]["error"]


def test_stale_processing_batch_is_swept_to_error(make_client, seeded_space, db):
    """ハートビートが停止した processing ジョブは取得時に error へ倒される。"""
    batch_id = "batch_stale"
    stale_at = (datetime.now(UTC) - timedelta(seconds=3600)).isoformat()
    db.document(f"spaces/{seeded_space}/integration_jobs/{batch_id}").set(
        {
            "batch_id": batch_id,
            "status": "processing",
            "stage": "interpret",
            "heartbeat_at": stale_at,
            "filenames": ["a.csv"],
        }
    )
    client = make_client(uid="uid_member")
    body = client.get(
        f"/api/integration/batches/{batch_id}", headers={"X-Space-Id": seeded_space}
    ).json()
    assert body["status"] == "error"
    assert "実行途絶" in body["error"]


def test_unknown_batch_returns_404(make_client, seeded_space):
    client = make_client(uid="uid_member")
    res = client.get("/api/integration/batches/batch_none", headers={"X-Space-Id": seeded_space})
    assert res.status_code == 404


def test_confirm_intent_classifies_free_text(make_client, seeded_space, monkeypatch):
    """テキスト入力（ボタンを押さない承認/取消）が意図判定エンドポイント経由で分類される。"""
    import routers.integration as integration_router

    calls: list[tuple[str, str]] = []

    async def _fake_classify(space, message, context_summary):
        calls.append((message, context_summary))
        return "approve" if "OK" in message else "cancel"

    monkeypatch.setattr(integration_router, "classify_confirmation_intent", _fake_classify)

    client = make_client(uid="uid_member")
    res = client.post(
        "/api/integration/confirm-intent",
        headers={"X-Space-Id": seeded_space},
        json={"message": "OK", "context_summary": "「展示会X」の取り込みプラン"},
    )
    assert res.status_code == 200
    assert res.json() == {"intent": "approve"}
    assert calls == [("OK", "「展示会X」の取り込みプラン")]

    res2 = client.post(
        "/api/integration/confirm-intent",
        headers={"X-Space-Id": seeded_space},
        json={"message": "やめて", "context_summary": "「展示会X」の取り込みプラン"},
    )
    assert res2.json() == {"intent": "cancel"}
