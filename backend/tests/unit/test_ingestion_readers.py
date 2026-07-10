"""ingestion.readers の単体テスト（ファイル形式判定 + Read ステージの抽出ロジック）。"""

import io

import docx as docx_lib
import pandas as pd
import pptx as pptx_lib
import pytest
from pptx.util import Inches

from ingestion import readers


def _make_docx(
    paragraphs: tuple[str, ...] = (), table_rows: list[list[str]] | None = None
) -> bytes:
    document = docx_lib.Document()
    for p in paragraphs:
        document.add_paragraph(p)
    if table_rows:
        table = document.add_table(rows=0, cols=len(table_rows[0]))
        for row in table_rows:
            cells = table.add_row().cells
            for i, value in enumerate(row):
                cells[i].text = value
    buf = io.BytesIO()
    document.save(buf)
    return buf.getvalue()


def _make_xlsx(rows: list[dict]) -> bytes:
    buf = io.BytesIO()
    pd.DataFrame(rows).to_excel(buf, index=False)
    return buf.getvalue()


def _make_pptx(texts: tuple[str, ...] = (), table_rows: list[list[str]] | None = None) -> bytes:
    presentation = pptx_lib.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    for text in texts:
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        box.text_frame.text = text
    if table_rows:
        rows, cols = len(table_rows), len(table_rows[0])
        table = slide.shapes.add_table(rows, cols, Inches(1), Inches(3), Inches(4), Inches(2)).table
        for r, row in enumerate(table_rows):
            for c, value in enumerate(row):
                table.cell(r, c).text = value
    buf = io.BytesIO()
    presentation.save(buf)
    return buf.getvalue()


def _make_pdf(page_texts: list[str]) -> bytes:
    """依存を増やさず `read_pdf` をテストするための最小 PDF（複数ページ、Base14 Helvetica のみ）。"""
    header = b"%PDF-1.4\n"
    page_obj_ids = list(range(3, 3 + len(page_texts)))
    font_obj_id = 3 + len(page_texts)
    content_obj_ids = list(range(font_obj_id + 1, font_obj_id + 1 + len(page_texts)))

    objects: dict[int, bytes] = {
        1: b"<< /Type /Catalog /Pages 2 0 R >>",
        2: f"<< /Type /Pages /Kids [{' '.join(f'{o} 0 R' for o in page_obj_ids)}] "
        f"/Count {len(page_texts)} >>".encode(),
        font_obj_id: b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    }
    for oid, content_oid in zip(page_obj_ids, content_obj_ids, strict=True):
        objects[oid] = (
            f"<< /Type /Page /Parent 2 0 R /Resources << /Font << /F1 {font_obj_id} 0 R >> >> "
            f"/MediaBox [0 0 612 792] /Contents {content_oid} 0 R >>"
        ).encode()
    for content_oid, text in zip(content_obj_ids, page_texts, strict=True):
        stream = f"BT /F1 12 Tf 72 700 Td ({text}) Tj ET".encode("latin-1")
        objects[content_oid] = (
            f"<< /Length {len(stream)} >>\nstream\n".encode() + stream + b"\nendstream"
        )

    body = b""
    offsets: dict[int, int] = {}
    for oid in sorted(objects):
        offsets[oid] = len(header) + len(body)
        body += f"{oid} 0 obj\n".encode() + objects[oid] + b"\nendobj\n"

    max_oid = max(objects)
    xref_start = len(header) + len(body)
    xref_lines = [f"0 {max_oid + 1}", "0000000000 65535 f "]
    xref_lines += [f"{offsets.get(oid, 0):010d} 00000 n " for oid in range(1, max_oid + 1)]
    xref = "xref\n" + "\n".join(xref_lines) + "\n"
    trailer = f"trailer\n<< /Size {max_oid + 1} /Root 1 0 R >>\nstartxref\n{xref_start}\n%%EOF"
    return header + body + xref.encode() + trailer.encode()


def test_is_supported_accepts_docx():
    assert readers.is_supported("overview.docx")
    assert readers.is_docx("overview.docx")
    assert not readers.is_tabular("overview.docx")


def test_ensure_supported_rejects_legacy_doc():
    with pytest.raises(readers.UnsupportedFileError):
        readers.ensure_supported("legacy.doc")


def test_is_supported_accepts_pdf_and_pptx():
    for name in ("report.pdf", "slides.pptx"):
        assert readers.is_supported(name)
        assert not readers.is_tabular(name)
    assert readers.is_pdf("report.pdf")
    assert not readers.is_pdf("slides.pptx")
    assert readers.is_pptx("slides.pptx")
    assert not readers.is_pptx("report.pdf")


def test_extraction_caveat_only_for_pdf_and_pptx():
    assert readers.extraction_caveat("report.pdf") != ""
    assert readers.extraction_caveat("slides.pptx") != ""
    for name in ("overview.docx", "list.xlsx", "notes.txt", "list.csv"):
        assert readers.extraction_caveat(name) == ""


def test_read_docx_extracts_paragraphs_and_tables():
    content = _make_docx(
        paragraphs=("概要テキスト",),
        table_rows=[["日程", "会場"], ["10/1", "東京"]],
    )
    text = readers.read_docx(content)
    assert "概要テキスト" in text
    assert "日程 | 会場" in text
    assert "10/1 | 東京" in text


def test_read_blocks_docx_returns_single_text_block():
    content = _make_docx(paragraphs=("概要",))
    blocks = readers.read_blocks("overview.docx", content)
    assert len(blocks) == 1
    assert blocks[0].row_no == 0
    assert "概要" in blocks[0].raw["text"]
    assert blocks[0].read_error == ""


def test_read_blocks_corrupt_docx_yields_read_error():
    blocks = readers.read_blocks("broken.docx", b"not a real docx file")
    assert len(blocks) == 1
    assert blocks[0].read_error.startswith("読み込みエラー")


def test_read_blocks_xlsx_returns_row_blocks():
    content = _make_xlsx(
        [{"会社名": "株式会社A", "お名前": "山田様"}, {"会社名": "株式会社B", "お名前": "鈴木様"}]
    )
    blocks = readers.read_blocks("list.xlsx", content)
    assert len(blocks) == 2
    assert blocks[0].read_error == ""
    assert blocks[0].raw == {"会社名": "株式会社A", "お名前": "山田様"}
    assert blocks[1].raw == {"会社名": "株式会社B", "お名前": "鈴木様"}


def test_read_blocks_corrupt_xlsx_yields_read_error():
    blocks = readers.read_blocks("broken.xlsx", b"not a real xlsx file")
    assert len(blocks) == 1
    assert blocks[0].read_error.startswith("読み込みエラー")


def test_read_pptx_extracts_text_and_tables():
    content = _make_pptx(
        texts=("概要テキスト",),
        table_rows=[["日程", "会場"], ["10/1", "東京"]],
    )
    text = readers.read_pptx(content)
    assert "概要テキスト" in text
    assert "日程 | 会場" in text
    assert "10/1 | 東京" in text


def test_read_blocks_pptx_returns_single_text_block():
    content = _make_pptx(texts=("概要",))
    blocks = readers.read_blocks("overview.pptx", content)
    assert len(blocks) == 1
    assert blocks[0].row_no == 0
    assert "概要" in blocks[0].raw["text"]
    assert blocks[0].read_error == ""


def test_read_blocks_corrupt_pptx_yields_read_error():
    blocks = readers.read_blocks("broken.pptx", b"not a real pptx file")
    assert len(blocks) == 1
    assert blocks[0].read_error.startswith("読み込みエラー")


def test_read_pdf_extracts_text_across_pages():
    content = _make_pdf(["Hello PDF", "Second page"])
    text = readers.read_pdf(content)
    assert "Hello PDF" in text
    assert "Second page" in text


def test_read_blocks_pdf_returns_single_text_block():
    content = _make_pdf(["Hello PDF"])
    blocks = readers.read_blocks("report.pdf", content)
    assert len(blocks) == 1
    assert blocks[0].row_no == 0
    assert "Hello PDF" in blocks[0].raw["text"]
    assert blocks[0].read_error == ""


def test_read_blocks_corrupt_pdf_yields_read_error():
    blocks = readers.read_blocks("broken.pdf", b"not a real pdf file")
    assert len(blocks) == 1
    assert blocks[0].read_error.startswith("読み込みエラー")
