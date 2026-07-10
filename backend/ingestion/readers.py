"""
readers — ファイル → 観測ブロック（Read ステージの読み込み部。純粋・Firestore 非依存）

対応形式は CSV / Excel / テキスト / Word (.docx) / PDF / PowerPoint (.pptx)。
PDF・PPTX はテキスト抽出のみ（表のレイアウト・図解・画像内の文字は取り込まれない場合が
ある。multimodal 読み取りは需要が実際に発生した時点の拡張トリガーのまま、ADR-015）。
抽出品質の注意は `extraction_caveat()` が返し、確認画面（Confirm）に表示する。
旧形式の .doc 等の未対応形式は UnsupportedFileError を送出し、ルーターが 400 に変換する。
"""

import io
from dataclasses import dataclass

import docx
import pandas as pd
import pptx
from pypdf import PdfReader

SUPPORTED_TABULAR = (".csv", ".xlsx", ".xls")
SUPPORTED_TEXT = (".txt",)
SUPPORTED_DOCX = (".docx",)
SUPPORTED_PDF = (".pdf",)
SUPPORTED_PPTX = (".pptx",)

_EXTRACTION_CAVEATS = {
    SUPPORTED_PDF: (
        "PDF はテキスト抽出のみのため、表のレイアウトや画像内の文字は取り込まれない"
        "場合があります。抽出結果をご確認ください。"
    ),
    SUPPORTED_PPTX: (
        "PowerPoint はテキスト抽出のみのため、図解・グラフ・画像内の文字は取り込まれない"
        "場合があります。抽出結果をご確認ください。"
    ),
}


class UnsupportedFileError(ValueError):
    """未対応のファイル形式（.doc 等）。ルーターが 400 に変換する。"""


def is_tabular(filename: str) -> bool:
    return filename.lower().endswith(SUPPORTED_TABULAR)


def is_docx(filename: str) -> bool:
    return filename.lower().endswith(SUPPORTED_DOCX)


def is_pdf(filename: str) -> bool:
    return filename.lower().endswith(SUPPORTED_PDF)


def is_pptx(filename: str) -> bool:
    return filename.lower().endswith(SUPPORTED_PPTX)


def is_supported(filename: str) -> bool:
    return filename.lower().endswith(
        SUPPORTED_TABULAR + SUPPORTED_TEXT + SUPPORTED_DOCX + SUPPORTED_PDF + SUPPORTED_PPTX
    )


def ensure_supported(filename: str) -> None:
    if not is_supported(filename):
        raise UnsupportedFileError(
            f"未対応のファイル形式です: {filename}"
            "（対応形式: CSV / Excel / テキスト / Word / PDF / PowerPoint）"
        )


def extraction_caveat(filename: str) -> str:
    """フォーマット起因の定型注意（PDF/PPTX のみ非空）。P1 の決定論的判定で、AI は関与しない。"""
    lower = filename.lower()
    for exts, message in _EXTRACTION_CAVEATS.items():
        if lower.endswith(exts):
            return message
    return ""


def read_tabular(filename: str, content: bytes) -> tuple[list[str], list[dict]]:
    """表形式ファイルを (ヘッダー, 行データ) で返す。値はすべて文字列。"""
    if filename.lower().endswith(".csv"):
        df = pd.read_csv(io.BytesIO(content), encoding="utf-8-sig", dtype=str)
    else:
        df = pd.read_excel(io.BytesIO(content), dtype=str)
    df = df.fillna("")
    return list(df.columns), df.to_dict(orient="records")


def read_text(content: bytes) -> str:
    return content.decode("utf-8", errors="replace")


def read_docx(content: bytes) -> str:
    """Word 文書 (.docx) から本文テキストを抽出する（段落 + 表）。

    表の各行は段落群の後に "セルA | セルB" 形式で連結する。ドキュメント順の厳密な
    インターリーブはしない（Interpret ステージは全文を1回の AI 呼び出しに渡すだけ
    のため、順序の厳密性より実装の単純さを優先した。品質が問題になれば
    iter_block_items での書き直しを検討する）。
    """
    document = docx.Document(io.BytesIO(content))
    parts: list[str] = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def read_pdf(content: bytes) -> str:
    """PDF からページ単位のテキストを抽出する（テキスト抽出のみ。表・レイアウトは崩れる前提）。"""
    reader = PdfReader(io.BytesIO(content))
    parts = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(p for p in parts if p.strip())


def read_pptx(content: bytes) -> str:
    """PowerPoint (.pptx) からテキストを抽出する（テキストボックス + 表）。

    docx の `read_docx` と同じ「段落＋表を "セルA | セルB" 形式で連結」の方針を踏襲する。
    スライド順に走査するが、スライド区切りのマーカーは付けない（実装の単純さを優先）。
    """
    presentation = pptx.Presentation(io.BytesIO(content))
    parts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
            elif shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
    return "\n".join(parts)


def read_document_text(filename: str, content: bytes) -> str:
    """非表形式ファイルから本文テキストを抽出する（.txt はデコード、.docx/.pdf/.pptx は専用抽出）。"""
    if is_docx(filename):
        return read_docx(content)
    if is_pdf(filename):
        return read_pdf(content)
    if is_pptx(filename):
        return read_pptx(content)
    return read_text(content)


@dataclass
class SourceBlock:
    """観測ブロック（source_records に着地する単位）。表=1行、文書=1件。"""

    row_no: int
    raw: dict  # {元列: 値}（文書は {"text": 全文}）
    read_error: str = ""  # 読み込み失敗の理由（非空なら skipped として着地）


def read_blocks(filename: str, content: bytes) -> list[SourceBlock]:
    """ファイルを観測ブロック列に変換する。未対応形式は UnsupportedFileError。"""
    ensure_supported(filename)
    if is_tabular(filename):
        try:
            _, rows = read_tabular(filename, content)
        except Exception as e:  # 壊れた表ファイル: 1ブロックの読み込みエラーとして着地
            return [SourceBlock(row_no=0, raw={}, read_error=f"読み込みエラー: {e}")]
        return [SourceBlock(row_no=i, raw=row) for i, row in enumerate(rows)]
    try:
        text = read_document_text(filename, content)
    except Exception as e:  # 壊れた docx 等: 1ブロックの読み込みエラーとして着地
        return [SourceBlock(row_no=0, raw={}, read_error=f"読み込みエラー: {e}")]
    return [SourceBlock(row_no=0, raw={"text": text})]
