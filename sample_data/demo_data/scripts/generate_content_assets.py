"""
content_catalog.txt に列挙された9件のコンテンツ資産を、それぞれ独立した個別ファイル
(ホワイトペーパー4件・導入事例3件 = PDF / セミナー・イベント案内2件 = PPTX)として生成する
ワンオフスクリプト。

「きれいに一覧化された資料は現場に存在せず、個々のコンテンツがバラバラに散在している」という
リアルなデモストーリーのための素材を作る。content_catalog.txt / seed_demo_space.py は変更しない。

PDF は日本語CID内蔵フォント(追加フォントファイル不要)の reportlab で生成する。reportlab は
デモデータ生成専用の一時依存なので backend 本番依存には入れず、実行時に --with で解決する。

実行:
    cd backend && uv run --with reportlab python \\
        ../sample_data/demo_data/scripts/generate_content_assets.py
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib.enums import TA_CENTER, TA_RIGHT
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.units import mm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.platypus import (
    HRFlowable,
    ListFlowable,
    ListItem,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
)

_OUT_DIR = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
    "content_assets",
)

# reportlab 内蔵の日本語CIDフォント(明朝=本文、ゴシック=見出し)。外部フォントファイル不要。
_MINCHO = "HeiseiMin-W3"
_GOTHIC = "HeiseiKakuGo-W5"

_BRAND = "クラウドフォージ株式会社"


# --------------------------------------------------------------------------- #
# PDF ヘルパ(reportlab)
# --------------------------------------------------------------------------- #
def _register_fonts() -> None:
    pdfmetrics.registerFont(UnicodeCIDFont(_MINCHO))
    pdfmetrics.registerFont(UnicodeCIDFont(_GOTHIC))


def _pdf_styles() -> dict[str, ParagraphStyle]:
    return {
        "eyebrow": ParagraphStyle(
            "eyebrow",
            fontName=_GOTHIC,
            fontSize=9,
            textColor="#5B6B7B",
            leading=13,
            spaceAfter=2,
        ),
        "title": ParagraphStyle(
            "title",
            fontName=_GOTHIC,
            fontSize=19,
            leading=27,
            textColor="#12283A",
            spaceAfter=4,
        ),
        "subtitle": ParagraphStyle(
            "subtitle",
            fontName=_MINCHO,
            fontSize=11,
            leading=17,
            textColor="#33475B",
            spaceAfter=6,
        ),
        "h2": ParagraphStyle(
            "h2",
            fontName=_GOTHIC,
            fontSize=13,
            leading=19,
            textColor="#12283A",
            spaceBefore=12,
            spaceAfter=5,
        ),
        "body": ParagraphStyle(
            "body",
            fontName=_MINCHO,
            fontSize=10.5,
            leading=17,
            textColor="#1F2933",
            spaceAfter=6,
        ),
        "bullet": ParagraphStyle(
            "bullet",
            fontName=_MINCHO,
            fontSize=10.5,
            leading=16,
            textColor="#1F2933",
        ),
        "meta": ParagraphStyle(
            "meta",
            fontName=_MINCHO,
            fontSize=8.5,
            leading=12,
            textColor="#8492A6",
        ),
        "footer_id": ParagraphStyle(
            "footer_id",
            fontName=_MINCHO,
            fontSize=8,
            leading=11,
            textColor="#8492A6",
            alignment=TA_RIGHT,
        ),
        "kpi": ParagraphStyle(
            "kpi",
            fontName=_GOTHIC,
            fontSize=12,
            leading=18,
            textColor="#0B6E5A",
            alignment=TA_CENTER,
            spaceAfter=2,
        ),
    }


def _bullets(items: list[str], style: ParagraphStyle) -> ListFlowable:
    return ListFlowable(
        [ListItem(Paragraph(t, style), leftIndent=6) for t in items],
        bulletType="bullet",
        bulletColor="#0B6E5A",
        bulletFontSize=7,
        leftIndent=12,
        spaceBefore=1,
        spaceAfter=6,
    )


def _build_pdf(
    filename: str,
    *,
    eyebrow: str,
    title: str,
    subtitle: str,
    meta: str,
    sections: list[tuple[str, list]],
    asset_id: str,
    url: str,
) -> None:
    """sections は (見出し, [Flowable...]) のリスト。"""
    s = _pdf_styles()
    path = os.path.join(_OUT_DIR, filename)

    def _footer(canvas, doc) -> None:
        canvas.saveState()
        canvas.setFont(_MINCHO, 8)
        canvas.setFillColor("#8492A6")
        canvas.drawString(20 * mm, 12 * mm, f"© 2026 {_BRAND}")
        canvas.drawRightString(190 * mm, 12 * mm, f"{asset_id}  |  {url}")
        canvas.drawCentredString(105 * mm, 12 * mm, f"- {doc.page} -")
        canvas.restoreState()

    story: list = [
        Paragraph(eyebrow, s["eyebrow"]),
        Paragraph(title, s["title"]),
        HRFlowable(width="100%", thickness=1.2, color="#0B6E5A", spaceAfter=6),
        Paragraph(subtitle, s["subtitle"]),
        Paragraph(meta, s["meta"]),
        Spacer(1, 6 * mm),
    ]
    for heading, flowables in sections:
        story.append(Paragraph(heading, s["h2"]))
        story.extend(flowables)

    SimpleDocTemplate(
        path,
        pagesize=A4,
        leftMargin=20 * mm,
        rightMargin=20 * mm,
        topMargin=18 * mm,
        bottomMargin=20 * mm,
        title=title,
        author=_BRAND,
    ).build(story, onFirstPage=_footer, onLaterPages=_footer)


# --------------------------------------------------------------------------- #
# ホワイトペーパー(PDF x4)
# --------------------------------------------------------------------------- #
def wp_costpilot_intro() -> None:
    s = _pdf_styles()
    _build_pdf(
        "wp_costpilot_intro.pdf",
        eyebrow="WHITE PAPER ｜ FinOps・クラウドコスト最適化",
        title="CostPilot 導入ガイドブック",
        subtitle="〜クラウド費用を20%削減する最初の一歩〜",
        meta="発行: 2026年05月 / 対象: インフラエンジニア・情報システム部門",
        asset_id="wp_costpilot_intro",
        url="https://example.com/resources/wp_costpilot_intro",
        sections=[
            (
                "なぜ今クラウドコスト最適化なのか",
                [
                    Paragraph(
                        "マルチクラウド化とオートスケールの普及により、クラウド費用は「使った分だけ」から"
                        "「気づけば膨らむ固定費」へと変質しています。多くの組織で、実際には使われていない"
                        "リソースや過剰なスペックのインスタンスが月額コストの2〜3割を占めているのが実情です。",
                        s["body"],
                    ),
                    Paragraph(
                        "CostPilot は AWS / GCP / Azure のコストを横断的に可視化し、未使用・過剰リソースを"
                        "AIが自動検知して具体的な削減アクションを提案する FinOps SaaS です。",
                        s["body"],
                    ),
                ],
            ),
            (
                "主要な課題",
                [
                    _bullets(
                        [
                            "どのサービス・どのチームがコストを押し上げているのか把握できない",
                            "検証環境の消し忘れ、アタッチされていないディスクなど「見えないムダ」が放置される",
                            "コスト削減の当番が特定の担当者に属人化し、継続的な改善が回らない",
                        ],
                        s["bullet"],
                    ),
                ],
            ),
            (
                "CostPilot が提供する機能",
                [
                    _bullets(
                        [
                            "マルチクラウド横断のコスト可視化ダッシュボード(タグ・チーム・サービス単位)",
                            "未使用リソース・過剰スペックのAI自動検知と、月額削減見込み額の提示",
                            "削減提案のワンクリック起票と、実施後の効果トラッキング",
                            "予算アラートと異常コストの早期通知",
                        ],
                        s["bullet"],
                    ),
                ],
            ),
            (
                "導入ステップ",
                [
                    _bullets(
                        [
                            "STEP1: 読み取り専用ロールでクラウドアカウントを連携(最短30分)",
                            "STEP2: 初回スキャンで削減見込み額と優先度の高い項目を確認",
                            "STEP3: 提案を実施し、翌月のコストで効果を検証・レポート化",
                        ],
                        s["bullet"],
                    ),
                ],
            ),
            (
                "まとめ・お問い合わせ",
                [
                    Paragraph(
                        "CostPilot 導入企業では、平均して導入後3ヶ月でクラウド費用の約20%削減を実現しています。"
                        "まずは無料アセスメントで、貴社環境の削減余地を可視化しませんか。"
                        "お問い合わせは営業担当、または https://example.com/costpilot まで。",
                        s["body"],
                    ),
                ],
            ),
        ],
    )


def wp_observestack_mttr() -> None:
    s = _pdf_styles()
    _build_pdf(
        "wp_observestack_mttr.pdf",
        eyebrow="WHITE PAPER ｜ オブザーバビリティ・SRE",
        title="ObserveStack 活用ガイド",
        subtitle="〜MTTRを40%短縮した監視統合の実践〜",
        meta="発行: 2026年05月 / 対象: SREチーム・インフラ運用担当",
        asset_id="wp_observestack_mttr",
        url="https://example.com/resources/wp_observestack_mttr",
        sections=[
            (
                "分断された監視が障害対応を遅らせる",
                [
                    Paragraph(
                        "ログ・メトリクス・トレースがそれぞれ別々のツールに散らばっていると、障害発生時に"
                        "担当者は複数の画面を行き来しながら原因を推測することになります。この「ツール間の"
                        "行き来」こそが、MTTR(平均修復時間)を押し上げる最大の要因です。",
                        s["body"],
                    ),
                    Paragraph(
                        "ObserveStack は3つのシグナルを1つの画面に統合し、アラートから根本原因までを"
                        "1本の導線でたどれるようにする SRE 向け統合オブザーバビリティ SaaS です。",
                        s["body"],
                    ),
                ],
            ),
            (
                "主要な課題",
                [
                    _bullets(
                        [
                            "アラートは鳴るが、どのサービスのどの変更が原因か即座に切り分けられない",
                            "ログ検索・メトリクス確認・トレース追跡でツールを3つ跨ぐ",
                            "オンコール担当者ごとに調査の勘所が異なり、対応品質が安定しない",
                        ],
                        s["bullet"],
                    ),
                ],
            ),
            (
                "ObserveStack が提供する機能",
                [
                    _bullets(
                        [
                            "ログ・メトリクス・トレースの統合ビューと相関表示",
                            "アラートから関連トレース・ログへワンクリックで遷移する調査導線",
                            "サービスマップによる依存関係の可視化と影響範囲の即時把握",
                            "SLO / エラーバジェット管理とバーンレートアラート",
                        ],
                        s["bullet"],
                    ),
                ],
            ),
            (
                "導入ステップ",
                [
                    _bullets(
                        [
                            "STEP1: エージェント / OpenTelemetry でシグナルを集約",
                            "STEP2: 主要サービスのSLOとダッシュボードをテンプレートから設定",
                            "STEP3: 障害訓練(ゲームデー)で調査導線を検証し、Runbookに反映",
                        ],
                        s["bullet"],
                    ),
                ],
            ),
            (
                "まとめ・お問い合わせ",
                [
                    Paragraph(
                        "統合監視への移行により、ObserveStack 導入企業では障害対応時間(MTTR)を平均40%短縮"
                        "しています。詳しくは https://example.com/observestack をご覧ください。",
                        s["body"],
                    ),
                ],
            ),
        ],
    )


def wp_guardrail_iac() -> None:
    s = _pdf_styles()
    _build_pdf(
        "wp_guardrail_iac.pdf",
        eyebrow="WHITE PAPER ｜ IaCセキュリティ・コンプライアンス",
        title="GuardRail 機能紹介資料",
        subtitle="〜IaCセキュリティ自動チェックの仕組み〜",
        meta="発行: 2026年05月 / 対象: プラットフォーム・セキュリティ担当",
        asset_id="wp_guardrail_iac",
        url="https://example.com/resources/wp_guardrail_iac",
        sections=[
            (
                "設定ミスはデプロイ前に止める",
                [
                    Paragraph(
                        "クラウドのセキュリティインシデントの多くは、脆弱性ではなく「設定ミス」に起因します。"
                        "公開設定のままのストレージ、過剰なIAM権限、暗号化の無効化——これらは Terraform 等の"
                        "IaC コードの段階で静的に検知できれば、本番に到達する前に防げます。",
                        s["body"],
                    ),
                    Paragraph(
                        "GuardRail は IaC の静的解析によりセキュリティ設定ミスをデプロイ前に検知し、"
                        "監査対応の証跡も自動で残す IaC セキュリティ・コンプライアンス自動チェック SaaS です。",
                        s["body"],
                    ),
                ],
            ),
            (
                "主要な課題",
                [
                    _bullets(
                        [
                            "レビューでの目視チェックに依存し、設定ミスがすり抜ける",
                            "SOC2 / ISMS などの監査で、証跡の収集・整理に膨大な工数がかかる",
                            "ポリシーが文書化されず、チーム間で判断基準がばらつく",
                        ],
                        s["bullet"],
                    ),
                ],
            ),
            (
                "GuardRail が提供する機能",
                [
                    _bullets(
                        [
                            "Terraform / CloudFormation 等の静的解析による設定ミス検知",
                            "CI パイプラインへの組み込みと、Pull Request 上でのインライン指摘",
                            "業界標準ベンチマーク(CIS等)に沿ったポリシーテンプレート",
                            "監査対応向けのコンプライアンスレポート自動生成",
                        ],
                        s["bullet"],
                    ),
                ],
            ),
            (
                "導入ステップ",
                [
                    _bullets(
                        [
                            "STEP1: リポジトリと連携し、既存IaCをスキャンして現状を可視化",
                            "STEP2: 重大度に応じてブロッキング/警告のポリシーを設定",
                            "STEP3: CIに組み込み、監査レポートを定期出力する運用に移行",
                        ],
                        s["bullet"],
                    ),
                ],
            ),
            (
                "まとめ・お問い合わせ",
                [
                    Paragraph(
                        "GuardRail の導入により、規制業種のお客様で監査対応準備の工数を大幅に削減した事例が"
                        "あります。詳しくは https://example.com/guardrail をご覧ください。",
                        s["body"],
                    ),
                ],
            ),
        ],
    )


def wp_deployflow_dora() -> None:
    s = _pdf_styles()
    _build_pdf(
        "wp_deployflow_dora.pdf",
        eyebrow="WHITE PAPER ｜ CI/CD・デプロイ自動化",
        title="DeployFlow 完全ガイド",
        subtitle="〜DORAメトリクスで見るデプロイ頻度向上の実践〜",
        meta="発行: 2026年05月 / 対象: 開発リーダー・DevOpsエンジニア",
        asset_id="wp_deployflow_dora",
        url="https://example.com/resources/wp_deployflow_dora",
        sections=[
            (
                "デプロイの「怖さ」を取り除く",
                [
                    Paragraph(
                        "デプロイ頻度が上がらない組織の多くは、技術力の問題ではなく「リリースが怖い」という"
                        "心理的ハードルを抱えています。手戻りのリスクを下げる仕組み——カナリアリリースと"
                        "自動ロールバック——を整えることが、頻度向上の近道です。",
                        s["body"],
                    ),
                    Paragraph(
                        "DeployFlow は CI/CD パイプラインを高速化し、カナリアリリースと自動ロールバックで"
                        "安全にデプロイ頻度を高める CI/CD・デプロイ自動化 SaaS です。",
                        s["body"],
                    ),
                ],
            ),
            (
                "DORAメトリクスで現状を測る",
                [
                    _bullets(
                        [
                            "デプロイの頻度(Deployment Frequency)",
                            "変更のリードタイム(Lead Time for Changes)",
                            "変更失敗率(Change Failure Rate)",
                            "サービス復旧時間(Time to Restore Service)",
                        ],
                        s["bullet"],
                    ),
                ],
            ),
            (
                "DeployFlow が提供する機能",
                [
                    _bullets(
                        [
                            "パイプラインの並列化・キャッシュ最適化によるビルド高速化",
                            "トラフィックを段階的に移すカナリアリリース",
                            "指標悪化を検知した際の自動ロールバック",
                            "DORAメトリクスの自動計測とチーム別ダッシュボード",
                        ],
                        s["bullet"],
                    ),
                ],
            ),
            (
                "導入ステップ",
                [
                    _bullets(
                        [
                            "STEP1: 既存パイプラインを接続し、DORA4指標のベースラインを取得",
                            "STEP2: 主要サービスにカナリアリリースと自動ロールバックを設定",
                            "STEP3: 週次で指標を振り返り、ボトルネックを継続的に改善",
                        ],
                        s["bullet"],
                    ),
                ],
            ),
            (
                "まとめ・お問い合わせ",
                [
                    Paragraph(
                        "DeployFlow 導入企業では、安全性を保ったままデプロイ頻度を大きく向上させています。"
                        "詳しくは https://example.com/deployflow をご覧ください。",
                        s["body"],
                    ),
                ],
            ),
        ],
    )


# --------------------------------------------------------------------------- #
# 導入事例(PDF x3)
# --------------------------------------------------------------------------- #
def _case_kpis(kpis: list[tuple[str, str]], s: dict) -> list:
    """定量成果を「数字 + ラベル」で強調表示するブロック。"""
    out: list = []
    for value, label in kpis:
        out.append(Paragraph(value, s["kpi"]))
        out.append(
            Paragraph(
                label,
                ParagraphStyle("kpi_label", parent=s["meta"], alignment=TA_CENTER, spaceAfter=6),
            )
        )
    return out


def case_fintech_costpilot() -> None:
    s = _pdf_styles()
    _build_pdf(
        "case_fintech_costpilot.pdf",
        eyebrow="CASE STUDY ｜ CostPilot 導入事例",
        title="クラウド費用を月額240万円削減",
        subtitle="決済SaaSを提供するフィンテック企業様",
        meta="業種: フィンテック(決済SaaS) / 利用製品: CostPilot",
        asset_id="case_fintech_costpilot",
        url="https://example.com/cases/case_fintech_costpilot",
        sections=[
            (
                "成果サマリー",
                _case_kpis(
                    [
                        ("クラウド費用 −28%", "導入3ヶ月時点"),
                        ("月額 −240万円", "未使用リソースの自動検知による削減"),
                    ],
                    s,
                ),
            ),
            (
                "お客様概要",
                [
                    Paragraph(
                        "オンライン決済基盤をSaaSとして提供するフィンテック企業様。取引量の増加に合わせて"
                        "クラウドインフラを拡張してきた一方で、費用の増加ペースが売上を上回りつつあり、"
                        "コスト構造の見直しが経営課題となっていました。",
                        s["body"],
                    ),
                ],
            ),
            (
                "課題",
                [
                    _bullets(
                        [
                            "マルチアカウント構成で、どこにムダがあるか全体を俯瞰できない",
                            "スケール検証で立てた環境の消し忘れが常態化していた",
                            "コスト削減が特定エンジニアの手作業に依存し、継続しなかった",
                        ],
                        s["bullet"],
                    ),
                ],
            ),
            (
                "導入の決め手",
                [
                    Paragraph(
                        "読み取り専用ロールで短期間に導入でき、初回スキャンで削減見込み額が具体的な金額として"
                        "提示された点。「どこから手を付ければ効果が大きいか」が一目で分かり、経営層への説明も"
                        "容易になりました。",
                        s["body"],
                    ),
                ],
            ),
            (
                "成果",
                [
                    Paragraph(
                        "未使用リソースの自動検知と削減提案の実施により、導入3ヶ月でクラウド費用を28%(月額換算"
                        "で約240万円)削減。以降も予算アラートで異常コストを早期に抑えられる体制が定着しました。",
                        s["body"],
                    ),
                ],
            ),
            (
                "ご担当者コメント",
                [
                    Paragraph(
                        "「『節約』ではなく『可視化して意思決定する』文化に変わりました。削減額が数字で見えるので、"
                        "チームのモチベーションにもつながっています。」(インフラ責任者様)",
                        s["body"],
                    ),
                ],
            ),
        ],
    )


def case_ec_observestack() -> None:
    s = _pdf_styles()
    _build_pdf(
        "case_ec_observestack.pdf",
        eyebrow="CASE STUDY ｜ ObserveStack 導入事例",
        title="障害対応時間を平均45%短縮",
        subtitle="ECモール運営企業様",
        meta="業種: ECモール運営 / 利用製品: ObserveStack",
        asset_id="case_ec_observestack",
        url="https://example.com/cases/case_ec_observestack",
        sections=[
            (
                "成果サマリー",
                _case_kpis(
                    [
                        ("MTTR −45%", "障害対応時間の平均短縮"),
                        ("原因特定 大幅高速化", "ログ・メトリクス・トレース統合による"),
                    ],
                    s,
                ),
            ),
            (
                "お客様概要",
                [
                    Paragraph(
                        "複数ブランドが出店する大規模ECモールを運営する企業様。大型セール時にはアクセスが平常時の"
                        "数十倍に急増し、そのたびに一部サービスで障害が発生していました。",
                        s["body"],
                    ),
                ],
            ),
            (
                "課題",
                [
                    _bullets(
                        [
                            "セール時の障害で、原因がインフラかアプリか切り分けに時間がかかる",
                            "監視ツールが分散し、調査のたびに複数画面を行き来していた",
                            "オンコール担当者によって初動対応の質にばらつきがあった",
                        ],
                        s["bullet"],
                    ),
                ],
            ),
            (
                "導入の決め手",
                [
                    Paragraph(
                        "ログ・メトリクス・トレースを1画面に統合でき、アラートから関連トレースへワンクリックで"
                        "たどれる調査導線。セール本番という「待ったなし」の状況で効果を発揮すると判断しました。",
                        s["body"],
                    ),
                ],
            ),
            (
                "成果",
                [
                    Paragraph(
                        "統合ビューにより原因特定が高速化し、障害対応時間(MTTR)を平均45%短縮。次のセールでは"
                        "同種の障害を早期に封じ込め、機会損失の縮小につながりました。",
                        s["body"],
                    ),
                ],
            ),
            (
                "ご担当者コメント",
                [
                    Paragraph(
                        "「『どこを見ればいいか』で悩む時間がなくなりました。新しいメンバーでも同じ導線で調査"
                        "できるので、オンコールの心理的負担が下がっています。」(SREリード様)",
                        s["body"],
                    ),
                ],
            ),
        ],
    )


def case_sier_guardrail() -> None:
    s = _pdf_styles()
    _build_pdf(
        "case_sier_guardrail.pdf",
        eyebrow="CASE STUDY ｜ GuardRail 導入事例",
        title="監査対応準備工数を70%削減",
        subtitle="大手SIer子会社様",
        meta="業種: システムインテグレーション(金融・官公庁向け) / 利用製品: GuardRail",
        asset_id="case_sier_guardrail",
        url="https://example.com/cases/case_sier_guardrail",
        sections=[
            (
                "成果サマリー",
                _case_kpis(
                    [
                        ("監査準備工数 −70%", "SOC2 監査対応の準備"),
                        ("設定ミス デプロイ前に検知", "IaC静的解析の自動化による"),
                    ],
                    s,
                ),
            ),
            (
                "お客様概要",
                [
                    Paragraph(
                        "金融機関・官公庁向けのシステム開発を手がける大手SIerの子会社様。取り扱う案件の性質上、"
                        "高いセキュリティ水準と厳格な監査対応が求められていました。",
                        s["body"],
                    ),
                ],
            ),
            (
                "課題",
                [
                    _bullets(
                        [
                            "IaCの設定ミスをレビューの目視に頼っており、負荷が高い",
                            "SOC2監査のたびに証跡集めで多くの工数を消費していた",
                            "案件・チームごとにセキュリティ基準の解釈がばらついていた",
                        ],
                        s["bullet"],
                    ),
                ],
            ),
            (
                "導入の決め手",
                [
                    Paragraph(
                        "IaCの静的解析をCIに組み込み、設定ミスをデプロイ前に自動検知できる点。加えて、監査に"
                        "必要なコンプライアンスレポートを自動生成できることが、工数削減の決め手になりました。",
                        s["body"],
                    ),
                ],
            ),
            (
                "成果",
                [
                    Paragraph(
                        "セキュリティチェックの自動化とレポート自動生成により、SOC2監査対応の準備工数を70%削減。"
                        "属人化していた判断基準もポリシーとして明文化され、品質が安定しました。",
                        s["body"],
                    ),
                ],
            ),
            (
                "ご担当者コメント",
                [
                    Paragraph(
                        "「監査のたびに発生していた『証跡かき集め』がほぼなくなりました。開発者もPR上で指摘を"
                        "受けられるので、手戻りが減っています。」(セキュリティ管掌部門様)",
                        s["body"],
                    ),
                ],
            ),
        ],
    )


# --------------------------------------------------------------------------- #
# セミナー・イベント案内(PPTX x2)
# --------------------------------------------------------------------------- #
def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.text = bullets[0]
    for b in bullets[1:]:
        p = body.add_paragraph()
        p.text = b
        p.level = 0


def _add_footer_note(prs: Presentation, note: str) -> None:
    """最終スライド下部に asset_id / URL を控えめに記載する。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(9), Inches(1))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.text = note
    tf.paragraphs[0].font.size = Pt(11)


def seminar_deployflow_2026_08() -> None:
    prs = Presentation()
    _add_title_slide(
        prs,
        "【無料ウェビナー】DeployFlow入門セミナー",
        "〜デプロイ頻度を3倍にする自動化の実践〜 ｜ 2026年08月06日(木) 14:00-15:00",
    )
    _add_bullet_slide(
        prs,
        "開催概要",
        [
            "日時: 2026年08月06日(木) 14:00〜15:00(60分)",
            "形式: オンライン(Zoomウェビナー)",
            "参加費: 無料(事前登録制)",
            "対象: DevOpsエンジニア・開発リーダー",
            "主催: クラウドフォージ株式会社",
        ],
    )
    _add_bullet_slide(
        prs,
        "このセミナーで分かること",
        [
            "デプロイ頻度が上がらない組織に共通する「つまずき」",
            "カナリアリリースと自動ロールバックの基本的な仕組み",
            "DORAメトリクスを使った改善サイクルの回し方",
            "DeployFlow導入企業のビフォー・アフター",
        ],
    )
    _add_bullet_slide(
        prs,
        "アジェンダ",
        [
            "14:00 オープニング / DORAメトリクスの基礎",
            "14:15 デモ: カナリアリリースと自動ロールバック",
            "14:40 導入の進め方とよくある質問",
            "14:55 Q&A・クロージング",
        ],
    )
    _add_bullet_slide(
        prs,
        "お申し込み",
        [
            "参加費無料・以下より事前登録をお願いします",
            "申込URL: https://example.com/events/seminar_deployflow_2026_08",
            "お問い合わせ: マーケティング部 宮下",
        ],
    )
    _add_footer_note(
        prs,
        "asset_id: seminar_deployflow_2026_08\n"
        "URL: https://example.com/events/seminar_deployflow_2026_08\n"
        "コンテンツ種別: 未来のセミナー(募集中) / 開催予定日: 2026年08月06日(木)",
    )
    prs.save(os.path.join(_OUT_DIR, "seminar_deployflow_2026_08.pptx"))


def event_cloudops_autumn_2026() -> None:
    prs = Presentation()
    _add_title_slide(
        prs,
        "Cloud Ops Summit 2026 秋 出展のご案内",
        "クラウドフォージ ブースにお越しください ｜ 2026年10月13日(火)〜15日(木)",
    )
    _add_bullet_slide(
        prs,
        "開催概要",
        [
            "会期: 2026年10月13日(火)〜15日(木)",
            "種別: クラウド運用系カンファレンス(出展)",
            "出展社: クラウドフォージ株式会社",
            "対象: インフラエンジニア・SRE・情報システム部門・CTO/VPoE",
        ],
    )
    _add_bullet_slide(
        prs,
        "ブースで体験できること",
        [
            "CostPilot: クラウド費用の削減余地をその場で可視化",
            "ObserveStack: ログ・メトリクス・トレース統合のライブデモ",
            "GuardRail: IaC設定ミス検知のデモ",
            "DeployFlow: カナリアリリース・自動ロールバックのデモ",
        ],
    )
    _add_bullet_slide(
        prs,
        "見どころ",
        [
            "4製品の最新バージョンを一挙に展示",
            "エンジニア向けミニセッションを会期中に複数回実施",
            "個別相談ブースで貴社課題に合わせたデモをご用意",
        ],
    )
    _add_bullet_slide(
        prs,
        "来場のご案内",
        [
            "事前来場登録で当日スムーズにご入場いただけます",
            "詳細・登録URL: https://example.com/events/cloudops_summit_2026_autumn",
            "お問い合わせ: マーケティング部 岡田",
        ],
    )
    _add_footer_note(
        prs,
        "asset_id: event_cloudops_autumn_2026\n"
        "URL: https://example.com/events/cloudops_summit_2026_autumn\n"
        "コンテンツ種別: 未来のイベント(募集中) / 開催予定日: 2026年10月13日(火)〜15日(木)",
    )
    prs.save(os.path.join(_OUT_DIR, "event_cloudops_autumn_2026.pptx"))


# --------------------------------------------------------------------------- #
def main() -> None:
    os.makedirs(_OUT_DIR, exist_ok=True)
    _register_fonts()

    generators = [
        ("wp_costpilot_intro.pdf", wp_costpilot_intro),
        ("wp_observestack_mttr.pdf", wp_observestack_mttr),
        ("wp_guardrail_iac.pdf", wp_guardrail_iac),
        ("wp_deployflow_dora.pdf", wp_deployflow_dora),
        ("case_fintech_costpilot.pdf", case_fintech_costpilot),
        ("case_ec_observestack.pdf", case_ec_observestack),
        ("case_sier_guardrail.pdf", case_sier_guardrail),
        ("seminar_deployflow_2026_08.pptx", seminar_deployflow_2026_08),
        ("event_cloudops_autumn_2026.pptx", event_cloudops_autumn_2026),
    ]
    for _, fn in generators:
        fn()

    print(f"生成完了: {_OUT_DIR}")
    for name, _ in generators:
        path = os.path.join(_OUT_DIR, name)
        print(f"  {name}: {os.path.getsize(path)} bytes")


if __name__ == "__main__":
    main()
