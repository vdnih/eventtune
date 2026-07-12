"""
Cloud Ops Summit 2026 のデモ用リッチ形式ファイル(docx/xlsx/pptx)を生成するワンオフスクリプト。

backend の依存(python-docx / openpyxl / python-pptx)を再利用する。
実行: cd backend && uv run python ../demo_data/scripts/generate_rich_docs.py
"""

import os

from docx import Document
from openpyxl import Workbook
from openpyxl.styles import Font
from pptx import Presentation
from pptx.util import Inches, Pt

_OUT_DIR = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
    "01_cloud_ops_summit_2026",
)


def generate_overview_docx() -> None:
    doc = Document()
    doc.add_heading("Cloud Ops Summit 2026 展示会概要メモ", level=1)
    doc.add_paragraph("作成日: 2026年06月26日 / 担当: マーケティング部 岡田")

    doc.add_heading("イベント概要", level=2)
    for label, value in [
        ("イベント名", "Cloud Ops Summit 2026"),
        ("会場", "東京ビッグサイト 南展示棟"),
        ("ブース番号", "S-118"),
        ("開催期間", "2026年06月23日(火)〜25日(木)(3日間)"),
        ("イベント種別", "展示会"),
        ("開催状況", "終了"),
        ("予算", "4,200,000円"),
        ("目標集客数", "200名"),
    ]:
        p = doc.add_paragraph()
        p.add_run(f"{label}: ").bold = True
        p.add_run(value)

    doc.add_heading("出展目的", level=2)
    doc.add_paragraph(
        "CostPilot・ObserveStack・GuardRail・DeployFlowの4製品を横断的に訴求し、"
        "クラウド運用に課題を持つインフラエンジニア・SRE・情報システム部門・経営層との"
        "接点を作ることを目的とする。"
    )

    doc.add_heading("当日の所感・振り返りメモ", level=2)
    doc.add_paragraph(
        "3日間を通じて活気のある展示会だった。特にセキュリティゾーン(GuardRail)への反応が"
        "予想以上に良く、規制業種(金融・保険・SIer)からの「今すぐ導入を検討したい」という声を"
        "複数いただいた。コスト最適化ゾーン(CostPilot)への引き合いも初日から一定数あり、"
        "想定通りの反応だった。"
    )
    doc.add_paragraph(
        "2日目午後のデモ待ち行列が長くなりすぎた点は反省。次回はデモ担当者を最低3名に"
        "増員すべき。名刺回収は目標の150名に対し、来場者数の増加もあり終盤で対応しきれず"
        "取りこぼしが発生した可能性がある。"
    )
    doc.add_paragraph(
        "トーカイSIソリューションズ・保険テックラボの2社は特に決裁者本人が来場し、"
        "その場で見積もり依頼を受けるなど商談化の確度が高い。最優先でフォローする。"
    )

    doc.save(os.path.join(_OUT_DIR, "展示会概要メモ.docx"))


def generate_cost_xlsx() -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "費用実績"

    headers = ["費目", "内容", "金額(円)", "取引先", "請求日"]
    ws.append(headers)
    for cell in ws[1]:
        cell.font = Font(bold=True)

    rows = [
        ("会場費・出展費", "ブース出展料(S-118)", 2100000, "日本展示サービス株式会社", "2026-05-10"),
        ("ブース装飾・設営", "ブース装飾・什器設営一式", 950000, "株式会社ブースクラフト", "2026-06-20"),
        ("人件費・派遣", "マーケ部2名×3日間の人件費相当", 480000, "", ""),
        ("印刷・販促物・ノベルティ", "パンフレット300部・ノベルティTシャツ", 320000, "株式会社プリントワークス", "2026-06-15"),
        ("集客", "SNS広告・事前告知メール配信", 250000, "株式会社アドリーチ", "2026-06-05"),
        ("運営", "会場内Wi-Fiルーターレンタル等", 100000, "東京機材センター", "2026-06-20"),
    ]
    for row in rows:
        ws.append(row)

    # 合計行は表のデータ行として書かない(取り込み時に費用明細の1件として誤読され、
    # 合計額が二重計上される事故を避けるため)。別シートに参考値として記載する。
    summary_ws = wb.create_sheet("参考(合計)")
    summary_ws["A1"] = "費用合計(参考・取り込み対象外)"
    summary_ws["A2"] = sum(r[2] for r in rows)

    for col, width in zip("ABCDE", [22, 32, 14, 26, 14], strict=True):
        ws.column_dimensions[col].width = width

    wb.save(os.path.join(_OUT_DIR, "費用実績.xlsx"))


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def _add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.text = bullets[0]
    for b in bullets[1:]:
        p = body.add_paragraph()
        p.text = b
        p.level = 0


def generate_results_pptx() -> None:
    prs = Presentation()

    _add_title_slide(
        prs,
        "Cloud Ops Summit 2026 実施結果報告",
        "2026年06月26日 / マーケティング部",
    )

    _add_bullet_slide(
        prs,
        "KPIサマリー",
        [
            "ブース来場者数: 210名",
            "名刺獲得数: 45名(目標比 30%増)",
            "デモ実施数: 28件",
            "商談化見込み(アポ獲得): 14件",
            "アンケート回答: 38件 / NPS 8.1点",
        ],
    )

    _add_bullet_slide(
        prs,
        "興味関心の内訳(45名)",
        [
            "クラウドコスト最適化(CostPilot): 11名",
            "監視・障害対応(ObserveStack): 11名",
            "IaCセキュリティ(GuardRail): 10名",
            "CI/CD・デプロイ自動化(DeployFlow): 8名",
            "情報収集のみ: 5名",
        ],
    )

    _add_bullet_slide(
        prs,
        "Next Steps",
        [
            "温度感A(決裁権あり・高関心)を最優先でフォローメール送付",
            "トーカイSIソリューションズ・保険テックラボは見積もり提示を急ぐ",
            "セキュリティゾーンの反応が良く、次回は展示スペースを拡大検討",
            "7月のウェビナー・CTO円卓会議へ継続接点を作る",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    tf = tx.text_frame
    tf.text = "付録: ブース配置"
    tf.paragraphs[0].font.size = Pt(24)

    prs.save(os.path.join(_OUT_DIR, "実施結果報告.pptx"))


def main() -> None:
    os.makedirs(_OUT_DIR, exist_ok=True)
    generate_overview_docx()
    generate_cost_xlsx()
    generate_results_pptx()
    print(f"生成完了: {_OUT_DIR}")
    for name in ["展示会概要メモ.docx", "費用実績.xlsx", "実施結果報告.pptx"]:
        path = os.path.join(_OUT_DIR, name)
        print(f"  {name}: {os.path.getsize(path)} bytes")


if __name__ == "__main__":
    main()
